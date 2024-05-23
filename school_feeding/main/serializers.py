from django.contrib.auth.tokens import default_token_generator
from django.contrib.sites.shortcuts import get_current_site
from django.core.mail import send_mail
from django.template.loader import render_to_string
from django.utils.encoding import force_bytes
from django.utils.http import urlsafe_base64_encode
from rest_framework import serializers
from rest_framework.exceptions import ValidationError
from rest_framework.reverse import reverse
from django.contrib.auth.hashers import make_password
from pptx import Presentation


# changes
from .models import (  # PasswordReset,; PasswordResetConfirm,
    ApplicationForm,
    CustomUser,
    Message,
    PasswordChange,
    Project1,
    Project2,
    Project3,
    Project4,
    Registration,
    Team,
    VotableItem,
    Vote,
)

from PyPDF2 import PdfReader
from docx import Document
from io import BytesIO

class ForgetPasswordSerializer(serializers.Serializer):
    email = serializers.EmailField()

    def validate_email(self, value):
        user = CustomUser.objects.filter(email=value).first()
        if not user:
            raise serializers.ValidationError(
                "This email is not associated with any account."
            )
        return value

    # def send_reset_email(self, user, request):
    #     current_site = get_current_site(request)
    #     uid = urlsafe_base64_encode(force_bytes(user.pk))
    #     token = default_token_generator.make_token(user)
    #     reset_link = reverse("new_password", kwargs={"uidb64": uid, "token": token})
    #     reset_url = f"http://{current_site.domain}{reset_link}"
    #     subject = "Reset your password"
    #     message = f"Hello,\n\nPlease click on the following link to reset your password:\n\n{reset_url}\n\nIf you did not request a password reset, please ignore this email.\n\nThank you."
    #     send_mail(
    #         subject, message, "from@example.com", [user.email], fail_silently=False
    #     )
    #     print(
    #         f"Reset password link for user {user.username}: http://{current_site.domain}/reset_password/{uid}/{token}"
    #     )

    def send_reset_email(self, user, request):
        current_site = get_current_site(request)
        frontend_url = "https://aroxj_aprelakerpi_despan.schoolfeeding.am/newpass/"
        reset_link = f"{frontend_url}?uidb64={urlsafe_base64_encode(force_bytes(user.pk))}&token={default_token_generator.make_token(user)}"

        # reset_link = reverse("new_password", kwargs={"uidb64": uid, "token": token})
        subject = "Reset your password"
        message = f"Hello,\n\nPlease click on the following link to reset your password:\n\n{reset_link}\n\nIf you did not request a password reset, please ignore this email.\n\nThank you."
        from_email = "armanghandilyan977@gmail.com"
        send_mail(subject, message, from_email, [user.email], fail_silently=False)


class NewPasswordSerializer(serializers.Serializer):
    password = serializers.CharField(write_only=True)


class LoginSerializer(serializers.Serializer):
    email = serializers.EmailField()
    password = serializers.CharField()

    def validate(self, data):
        email = data.get("email")
        password = data.get("password")

        if email and password:
            # Authenticate user
            user = authenticate(email=email, password=password)

            if user:
                if not user.is_active:
                    raise serializers.ValidationError("User account is disabled.")
            else:
                raise serializers.ValidationError(
                    "Unable to log in with provided credentials."
                )
        else:
            raise serializers.ValidationError("Must provide both email and password.")

        data["user"] = user
        return data
        

class PasswordChangeSerializer(serializers.Serializer):

    class Meta:
        model = PasswordChange
        field = "__all__"

    password = serializers.CharField(max_length=110)
    new_password = serializers.CharField(max_length=110)


class ApplicationFormSerializer(serializers.ModelSerializer):

    class Meta:
        model = ApplicationForm
        fields = "__all__"




class UserSerializer(serializers.ModelSerializer):
    password = serializers.CharField(write_only=True)

    def create(self, validated_data):
        password = validated_data.pop("password")
        hashed_password = make_password(password)
        user = super().create({**validated_data, "password": hashed_password})
        return user

    class Meta:
        model = CustomUser
        fields = ["username", "email", "phone", "password"]


class TeamSerializer(serializers.ModelSerializer):

    class Meta:
        model = Team
        fields = "__all__"


class ProjectSerializer1(serializers.ModelSerializer):

    class Meta:
        model = Project1
        fields = "__all__"

    def validate_pdf(self, value):
        if value:
            try:
                pdf_reader = PdfReader(value)
                num_pages = len(pdf_reader.pages)
                if num_pages >= 2:
                    raise serializers.ValidationError("PDF file cannot have more than 15 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing PDF file: {e}")

        return value


    def validate_word(self, value):
        if value:
            try:
                doc = Document(value)
                num_pages = len(doc.paragraphs)
                if num_pages > 3:
                    raise serializers.ValidationError("Word file must have at most 3 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing Word file: {e}")

        return value

    def create(self, validated_data):
        pdf_file = validated_data.pop('pdf')
        word_file = validated_data.pop('word')


        project = Project1(pdf=pdf_file, word=word_file, **validated_data)
        project.save()
        return project


class ProjectSerializer2(serializers.ModelSerializer):

    class Meta:
        model = Project2
        fields = "__all__"

    
    def validate_pdf(self, value):
        if value:
            try:
                pdf_reader = PdfReader(value)
                num_pages = len(pdf_reader.pages)
                if num_pages >= 2:
                    raise serializers.ValidationError("PDF file cannot have more than 15 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing PDF file: {e}")

        return value


    def validate_word(self, value):
        if value:
            try:
                doc = Document(value)
                num_pages = len(doc.paragraphs)
                if num_pages > 3:
                    raise serializers.ValidationError("Word file must have at most 3 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing Word file: {e}")

        return value

    def create(self, validated_data):
        pdf_file = validated_data.pop('pdf')
        word_file = validated_data.pop('word')

        project = Project2(pdf=pdf_file, word=word_file, **validated_data)
        project.save()
        return project



class ProjectSerializer3(serializers.ModelSerializer):

    class Meta:
        model = Project3
        fields = "__all__"

    
    def validate_pdf(self, value):
        if value:
            try:
                # Check if the file is a PDF
                if value.name.endswith('.pdf'):
                    pdf_reader = PdfReader(value)
                    num_pages = len(pdf_reader.pages)
                    if num_pages >= 15:
                        raise serializers.ValidationError("PDF file cannot have more than 15 pages.")
                # Check if the file is a PowerPoint presentation
                elif value.name.endswith('.ppt') or value.name.endswith('.pptx'):
                    prs = Presentation(value)
                    num_slides = len(prs.slides)
                    if num_slides >= 15:
                        raise serializers.ValidationError("PowerPoint file must have at most 10 slides.")
                else:
                    raise serializers.ValidationError("Unsupported file format.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing file: {e}")

        return value

    def validate_pdf_only(self, value):
        if value:
            try:
                pdf_reader = PdfReader(value)
                num_pages = len(pdf_reader.pages)
                if num_pages >= 2:
                    raise serializers.ValidationError("PDF file cannot have more than 15 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing PDF file: {e}")

        return value

    def validate_word(self, value):
        if value:
            try:
                doc = Document(value)
                num_pages = len(doc.paragraphs)
                if num_pages > 3:
                    raise serializers.ValidationError("Word file must have at most 3 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing Word file: {e}")

        return value

    def create(self, validated_data):
        pdf_file = validated_data.pop('pdf')
        pdf_only_file = validated_data.pop('pdf_only')
        word_file = validated_data.pop('word')

        project = Project3(pdf=pdf_file, pdf_only=pdf_only_file, word=word_file, **validated_data)
        project.save()
        return project


class ProjectSerializer4(serializers.ModelSerializer):

    class Meta:
        model = Project4
        fields = "__all__"

    
    def validate_pdf(self, value):
        if value:
            try:
                # Check if the file is a PDF
                if value.name.endswith('.pdf'):
                    pdf_reader = PdfReader(value)
                    num_pages = len(pdf_reader.pages)
                    if num_pages >= 10:
                        raise serializers.ValidationError("PDF file cannot have more than 15 pages.")
                # Check if the file is a PowerPoint presentation
                elif value.name.endswith('.ppt') or value.name.endswith('.pptx'):
                    prs = Presentation(value)
                    num_slides = len(prs.slides)
                    if num_slides >= 10:
                        raise serializers.ValidationError("PowerPoint file must have at most 10 slides.")
                else:
                    raise serializers.ValidationError("Unsupported file format.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing file: {e}")

        return value

    def validate_pdf_only(self, value):
        if value:
            try:
                pdf_reader = PdfReader(value)
                num_pages = len(pdf_reader.pages)
                if num_pages >= 2:
                    raise serializers.ValidationError("PDF file cannot have more than 15 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing PDF file: {e}")

        return value

    def validate_word(self, value):
        if value:
            try:
                doc = Document(value)
                num_pages = len(doc.paragraphs)
                if num_pages > 3:
                    raise serializers.ValidationError("Word file must have at most 3 pages.")
            except Exception as e:
                raise serializers.ValidationError(f"Error processing Word file: {e}")

        return value

    def create(self, validated_data):
        pdf_file = validated_data.pop('pdf')
        pdf_only_file = validated_data.pop('pdf_only')
        word_file = validated_data.pop('word')

        project = Project4(pdf=pdf_file, pdf_only=pdf_only_file, word=word_file, **validated_data)
        project.save()
        return project


class RegistrationSerializer(serializers.ModelSerializer):

    class Meta:
        model = Registration
        fields = "__all__"


class MessageSerializer(serializers.ModelSerializer):

    class Meta:
        model = Message
        fields = "__all__"


# changes


class VotableItemSerializer(serializers.ModelSerializer):
    class Meta:
        model = VotableItem
        fields = "__all__"


class VoteSerializer(serializers.ModelSerializer):
    class Meta:
        model = Vote
        fields = "__all__"
